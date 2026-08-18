"""
Document Parser Service — Trích xuất text từ PDF/DOCX/PPTX và chunking.

Luồng xử lý:
    Upload file (PDF/DOCX/PPTX)
        ↓
    Lưu file trên disk (uploads/)
        ↓
    DB lưu metadata (filename, doc_type, file_path, status)
        ↓
    DocumentParser.parse(document) -> text
        ↓
    chunk_text(text) -> list[str]  (~1000 tokens/chunk, overlap 150)
        ↓
    Lưu chunks vào assessments.chunks (JSONB)
        ↓
    User chọn Persona → AI đọc chunks + prompt → sinh 10 câu hỏi

Tham khảo:
    PyPDF2:       https://pypdf2.readthedocs.io
    python-docx:  https://python-docx.readthedocs.io
    python-pptx:  https://python-pptx.readthedocs.io
"""
from io import BytesIO

import logging
import zipfile
from dataclasses import dataclass, field
from typing import List, Optional

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from app.services.vision_read import ImagePart, ReadResult

from app.models.document import DocType, Document
from app.services.storage import get_doc
from app.services.vision_read import read_file as vision_read_file

try:
    import rarfile
except ImportError:  # pragma: no cover
    rarfile = None  # type: ignore

logger = logging.getLogger(__name__)


@dataclass
class ParseResult:
    """Kết quả parse: text thường + mô tả diagram (từ vision reader)."""
    text: str
    diagrams: List[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
# Số ký tự ước lượng cho ~1000 tokens (tiếng Anh ~4 chars/token, tiếng Việt ~2)
CHUNK_SIZE_CHARS = 4000      # ~1000 tokens
CHUNK_OVERLAP_CHARS = 600    # ~150 tokens overlap để giữ ngữ cảnh

# Ngưỡng cảnh báo khi trích xuất được quá ít text (file scan, file rỗng, ...)
MIN_TEXT_LENGTH_WARN = 50


class DocumentParserError(Exception):
    """Lỗi khi parse document (file không đọc được, format sai, ...)."""


# ---------------------------------------------------------------------------
# Extractors — mỗi loại file một hàm riêng
# ---------------------------------------------------------------------------
def _extract_pdf(src) -> str:
    """Trích xuất text từ PDF. Nhận str path hoặc file-like object."""
    reader = PdfReader(src)
    pages: List[str] = []
    for idx, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # một trang lỗi không chặn cả file
            logger.warning("PDF page %s extract failed: %s", idx, exc)
            text = ""
        pages.append(text)
    return "\n\n".join(pages).strip()


def _extract_docx(src) -> str:
    """Trích xuất text từ DOCX. Nhận str path hoặc file-like object."""
    doc = DocxDocument(src)

    parts: List[str] = [p.text for p in doc.paragraphs if p.text]

    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts).strip()


def _extract_pptx(src) -> str:
    """Trích xuất text từ PPTX. Nhận str path hoặc file-like object."""
    prs = Presentation(src)
    blocks: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_lines: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
        if slide_lines:
            blocks.append(f"### Slide {slide_idx}\n" + "\n".join(slide_lines))
    return "\n\n".join(blocks).strip()


def _extract_md(src) -> str:
    """Trích xuất text từ Markdown. Plain text, bỏ qua markup."""
    import io
    content = src.read() if isinstance(src, io.IOBase) else open(src, "rb").read()
    raw = content.decode("utf-8", errors="replace")
    return raw.strip()


# File extensions đọc được text trong archive (code + text + office docs)
TEXT_EXTENSIONS = {
    # code
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rb", ".php", ".cs",
    ".cpp", ".c", ".h", ".hpp", ".html", ".css", ".json", ".yml", ".yaml",
    ".xml", ".toml", ".ini", ".cfg", ".conf", ".sh", ".bat", ".ps1", ".sql",
    ".dockerfile", ".makefile", ".env", ".properties",
    # text / markdown
    ".md", ".markdown", ".txt", ".rst", ".log",
}
# Office docs trong archive: extract riêng theo từng loại
OFFICE_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
# Archive con lồng nhau (đọc đệ quy)
NESTED_ARCHIVE_EXTENSIONS = {".zip", ".rar"}

MAX_ARCHIVE_MEMBERS = 500
MAX_ARCHIVE_TOTAL_TEXT = 200 * 1024 * 1024  # 200MB text tổng


def _extract_zip(src) -> str:
    """Trích xuất text từ ZIP: đọc mọi file text/code/office bên trong rồi ghép lại.

    Mỗi member được đánh dấu bằng header "### <path>" để AI biết nội dung từ đâu.
    Nhận str path, bytes hoặc file-like object (như các extractor khác).
    """
    data = src if isinstance(src, bytes) else src.read()
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            infos = [i for i in archive.infolist() if not i.is_dir()]
            if len(infos) > MAX_ARCHIVE_MEMBERS:
                raise DocumentParserError(
                    f"ZIP chứa quá nhiều file ({len(infos)} > {MAX_ARCHIVE_MEMBERS})"
                )
            parts: List[str] = []
            total = 0
            for info in infos:
                name = info.filename
                ext = name.split(".")[-1].lower() if "." in name else ""
                ext = f".{ext}"
                total += info.file_size
                if total > MAX_ARCHIVE_TOTAL_TEXT:
                    break
                try:
                    raw = archive.read(info)
                except Exception:
                    continue  # member lỗi thì bỏ qua, không chặn cả file
                text = _extract_member_text(raw, name, ext)
                if text:
                    parts.append(f"### {name}\n{text}")
            return "\n\n".join(parts).strip()
    except zipfile.BadZipFile as exc:
        raise DocumentParserError("File ZIP bị lỗi hoặc không thể giải nén") from exc


def _extract_rar(data: bytes) -> str:
    """Trích xuất text từ RAR (nếu có rarfile), tương tự _extract_zip."""
    if rarfile is None:
        raise DocumentParserError("Chưa cài thư viện rarfile để đọc file .rar")
    try:
        with rarfile.RarFile(BytesIO(data)) as archive:
            infos = [i for i in archive.infolist() if not i.isdir()]
            if len(infos) > MAX_ARCHIVE_MEMBERS:
                raise DocumentParserError(
                    f"RAR chứa quá nhiều file ({len(infos)} > {MAX_ARCHIVE_MEMBERS})"
                )
            parts: List[str] = []
            total = 0
            for info in infos:
                name = info.filename
                ext = name.split(".")[-1].lower() if "." in name else ""
                ext = f".{ext}"
                total += info.file_size
                if total > MAX_ARCHIVE_TOTAL_TEXT:
                    break
                try:
                    raw = archive.read(name)
                except Exception:
                    continue
                text = _extract_member_text(raw, name, ext)
                if text:
                    parts.append(f"### {name}\n{text}")
            return "\n\n".join(parts).strip()
    except (rarfile.RarCannotExec, rarfile.Error) as exc:
        raise DocumentParserError(f"File RAR bị lỗi hoặc không thể giải nén: {exc}") from exc


def _extract_member_text(raw: bytes, name: str, ext: str) -> str:
    """Trích text 1 member trong archive theo extension."""
    try:
        if ext in OFFICE_EXTENSIONS:
            if ext == ".pdf":
                return _extract_pdf(BytesIO(raw))
            if ext == ".docx":
                return _extract_docx(BytesIO(raw))
            if ext == ".pptx":
                return _extract_pptx(BytesIO(raw))
            if ext == ".xlsx":
                return _extract_xlsx(BytesIO(raw))
        if ext in NESTED_ARCHIVE_EXTENSIONS:
            return _extract_zip(raw) if ext == ".zip" else _extract_rar(raw)
        if ext in TEXT_EXTENSIONS or ext == "":
            return raw.decode("utf-8", errors="replace").strip()
    except Exception:
        logger.debug("Skip member %s: %s", name, exc_info=True)
    return ""


def _extract_archive(src) -> str:
    """Nhận ZIP hoặc RAR, tự detect theo magic bytes và dispatch."""
    data = src if isinstance(src, bytes) else src.read()
    if data[:8].startswith(b"Rar!\x1a\x07"):
        return _extract_rar(data)
    return _extract_zip(data)


def _extract_xlsx(src) -> str:
    """Trích xuất text từ XLSX (openpyxl nếu có, không thì bỏ qua)."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(BytesIO(src.read()) if isinstance(src, BytesIO) else src, read_only=True, data_only=True)
    except Exception:
        return ""
    rows: List[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None and str(c).strip()]
            if vals:
                rows.append("\t".join(vals))
    return "\n".join(rows).strip()


_EXTRACTORS = {
    DocType.PDF: _extract_pdf,
    DocType.DOCX: _extract_docx,
    DocType.PPTX: _extract_pptx,
    DocType.ZIP: _extract_archive,
    # .md mapped to DocType.PDF via router, handle via filename fallback below
}


# ---------------------------------------------------------------------------
# Public API — dùng từ router / service khác
# ---------------------------------------------------------------------------
# Doc types routed PRIMARY through Gemini 3.1 Flash Lite vision reader.
# Native extractors (PyPDF2/python-docx) only read text — they drop diagrams,
# ERD/DFD/sequence shapes embedded as images. Vision reader captures both.
VISION_PRIMARY_TYPES = {DocType.PDF, DocType.DOCX, DocType.PPTX}

# Doc types that Gemini cannot meaningfully parse (archives, plain text).
# Keep native extractors as the only path for these.
_NATIVE_ONLY_TYPES = {DocType.ZIP}


async def extract_text(document) -> ParseResult:
    """
    Đọc file từ MinIO theo document.doc_type và trích xuất text + diagram.

    Stage 5 (Option B): pdf/docx/pptx → PRIMARY qua Gemini 3.1 Flash Lite vision reader.
    - PDF: gửi nguyên file, Gemini render ảnh nhúng trực tiếp (text + diagram 1 call).
    - DOCX/PPTX: Gemini KHÔNG render được ảnh nhúng trong OOXML khi gửi nguyên file
      (chỉ đọc document.xml). Nên unzip, trích ảnh raster, gửi TỪNG ảnh như 1 part
      multimodal kèm native text làm context → Gemini mô tả từng sơ đồ (Option B).
    - zip/md/text → giữ parser cũ (Gemini không đọc archive/code).

    Args:
        document: ORM Document (cần doc_type + storage_key).

    Returns:
        ParseResult(text, diagrams). diagrams rỗng nếu file không có sơ đồ.

    Raises:
        DocumentParserError: nếu doc_type không hỗ trợ hoặc file lỗi.
    """
    storage_key = document.storage_key
    if not storage_key:
        raise DocumentParserError(f"Document {document.id} has no storage key")

    try:
        data = await get_doc(storage_key)
    except Exception as exc:
        logger.exception("Failed to download %s from MinIO", storage_key)
        raise DocumentParserError(f"Download failed for {storage_key}: {exc}") from exc

    # Markdown / plain text: decode directly, no vision needed.
    if storage_key.lower().endswith(".md"):
        text = data.decode("utf-8", errors="replace").strip()
        if len(text) < MIN_TEXT_LENGTH_WARN:
            logger.warning("Extracted text is suspiciously short (%s chars) from %s", len(text), storage_key)
        return ParseResult(text=text, diagrams=[])

    # ── PRIMARY: pdf/docx/pptx → Gemini vision reader ──
    if document.doc_type in VISION_PRIMARY_TYPES:
        mime_type = _doc_type_to_mime(document.doc_type)
        if mime_type:
            if document.doc_type == DocType.PDF:
                # PDF: Gemini renders embedded images natively from the whole file.
                try:
                    vision = await vision_read_file(data, mime_type)
                    if vision.text or vision.diagrams:
                        logger.info(
                            "Vision reader extracted %d chars, %d diagrams from %s (PRIMARY path)",
                            len(vision.text), len(vision.diagrams), storage_key,
                        )
                        return ParseResult(text=vision.text, diagrams=vision.diagrams)
                    logger.warning(
                        "Vision reader returned empty for %s — falling back to native parser",
                        storage_key,
                    )
                except Exception as exc:
                    logger.warning(
                        "Vision reader failed for %s (%s) — falling back to native parser",
                        storage_key, exc,
                    )
            else:
                # DOCX/PPTX (Option B): Gemini cannot render OOXML-embedded images when
                # the whole file is sent (it only reads document.xml text). Unzip the
                # office package, extract the embedded raster images, and send each as a
                # separate multimodal part alongside the native text as context.
                # On any vision failure we keep the native text (already extracted) —
                # we do NOT fall through to the old parser, which would just re-read XML.
                native_text = _extract_office_text(data, document.doc_type)
                images = _extract_office_images(data, document.doc_type)
                if images:
                    try:
                        vision = await vision_read_file(
                            data, mime_type, images=images, body_text=native_text
                        )
                    except Exception as exc:
                        logger.warning(
                            "Vision reader failed for %s (%s) — using native text only",
                            storage_key, exc,
                        )
                        return ParseResult(text=native_text, diagrams=[])
                    if vision.diagrams:
                        logger.info(
                            "Vision reader described %d diagrams from %s (Option B, %d images)",
                            len(vision.diagrams), storage_key, len(images),
                        )
                        # text is the native text (authoritative); diagrams from vision.
                        return ParseResult(text=native_text, diagrams=vision.diagrams)
                    logger.warning(
                        "Vision reader found no diagrams in %d images of %s — "
                        "using native text only",
                        len(images), storage_key,
                    )
                else:
                    logger.info("No embedded images in %s — using native text only", storage_key)
                return ParseResult(text=native_text, diagrams=[])

    # ── FALLBACK / NATIVE-ONLY: parser cũ ──
    extractor = _EXTRACTORS.get(document.doc_type)
    if extractor is None:
        raise DocumentParserError(f"Unsupported doc_type: {document.doc_type}")

    file_io = BytesIO(data)
    try:
        text = extractor(file_io)
    except DocumentParserError:
        raise
    except Exception as exc:
        logger.exception("Failed to extract text from %s", storage_key)
        raise DocumentParserError(f"Extract failed for {storage_key}: {exc}") from exc

    if len(text) < MIN_TEXT_LENGTH_WARN:
        logger.warning(
            "Extracted text is suspiciously short (%s chars) from %s "
            "(file có thể là scan, ảnh, hoặc rỗng).",
            len(text), storage_key,
        )

    return ParseResult(text=text, diagrams=[])


def _doc_type_to_mime(doc_type: DocType) -> str | None:
    """Map DocType to MIME type for vision reader."""
    mapping = {
        DocType.PDF: "application/pdf",
        DocType.DOCX: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        DocType.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        DocType.ZIP: "application/zip",
    }
    return mapping.get(doc_type)


# ---------------------------------------------------------------------------
# Office image extraction (Option B) — Gemini cannot render OOXML-embedded images
# when the whole .docx/.pptx is sent, so we unzip and send each image separately.
# ---------------------------------------------------------------------------
# Media folders inside the OOXML zip.
_OFFICE_MEDIA_FOLDERS = {
    DocType.DOCX: "word/media",
    DocType.PPTX: "ppt/media",
}
# Max dimension (px) for an image sent to Gemini — keeps payloads small (Gemini
# inline_data limit is 20MB total; 82 images at full res would blow past it).
_MAX_IMAGE_DIM = 1600
# Skip tiny images (icons, bullets) that are not diagrams.
_MIN_IMAGE_DIM = 200


def _extract_office_text(data: bytes, doc_type: DocType) -> str:
    """Native text extraction for office files (used as vision context)."""
    extractor = _EXTRACTORS.get(doc_type)
    if extractor is None:
        return ""
    try:
        return extractor(BytesIO(data)) or ""
    except Exception as exc:  # native extractor failed — vision still works on images
        logger.warning("Native office text extraction failed for %s: %s", doc_type, exc)
        return ""


def _extract_office_images(data: bytes, doc_type: DocType) -> List[ImagePart]:
    """Unzip the office package and return embedded raster images as ImagePart list.

    Images are downscaled (longest side <= _MAX_IMAGE_DIM) and tiny images are skipped
    to stay within Gemini's inline_data limits and avoid sending icons/bullets.
    """
    folder = _OFFICE_MEDIA_FOLDERS.get(doc_type)
    if not folder:
        return []
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            parts: List[ImagePart] = []
            for info in archive.infolist():
                if info.is_dir() or not info.filename.startswith(folder + "/"):
                    continue
                ext = info.filename.rsplit(".", 1)[-1].lower()
                mime = {
                    "png": "image/png",
                    "jpg": "image/jpeg",
                    "jpeg": "image/jpeg",
                    "gif": "image/gif",
                    "webp": "image/webp",
                    "bmp": "image/bmp",
                }.get(ext)
                if not mime:
                    continue
                raw = archive.read(info)
                raw = _downscale_image(raw, mime)
                if raw is None:
                    continue
                parts.append(ImagePart(data=raw, mime_type=mime))
            logger.info("Extracted %d embedded images from %s", len(parts), doc_type)
            return parts
    except zipfile.BadZipFile as exc:
        logger.warning("Office file is not a valid zip (%s): %s", doc_type, exc)
        return []


def _downscale_image(raw: bytes, mime: str) -> Optional[bytes]:
    """Downscale image if too large; return (possibly resized) bytes or None if invalid."""
    try:
        from PIL import Image
    except ImportError:  # pragma: no cover — Pillow is a hard dep for Option B
        return raw
    try:
        with Image.open(BytesIO(raw)) as img:
            w, h = img.size
            if w < _MIN_IMAGE_DIM and h < _MIN_IMAGE_DIM:
                return None  # too small to be a diagram (icon/bullet)
            if max(w, h) > _MAX_IMAGE_DIM:
                scale = _MAX_IMAGE_DIM / max(w, h)
                img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
            out = BytesIO()
            save_fmt = "JPEG" if mime == "image/jpeg" else "PNG"
            img.convert("RGB" if save_fmt == "JPEG" else "RGBA").save(out, save_fmt)
            return out.getvalue()
    except Exception as exc:
        logger.debug("Skipping unreadable image (%s): %s", mime, exc)
        return None


def chunk_text(text: str,
               chunk_size: int = CHUNK_SIZE_CHARS,
               overlap: int = CHUNK_OVERLAP_CHARS) -> List[str]:
    """
    Cắt text dài thành các chunk ~1000 tokens (~4000 chars), overlap ~150 tokens.

    Cắt theo ranh giới tự nhiên (paragraph → sentence) để chunk gọn ý,
    tránh cắt giữa câu.

    Args:
        text: text đầu vào.
        chunk_size: số ký tự tối đa mỗi chunk (mặc định 4000 ≈ 1000 tokens).
        overlap: số ký tự overlap giữa 2 chunk liên tiếp (mặc định 600 ≈ 150).

    Returns:
        Danh sách string, mỗi phần tử là một chunk.
    """
    if not text or not text.strip():
        return []

    if chunk_size <= 0:
        raise ValueError("chunk_size must be > 0")
    if overlap < 0 or overlap >= chunk_size:
        raise ValueError("overlap must satisfy 0 <= overlap < chunk_size")

    # Đầu tiên tách theo paragraph, giữ lại paragraph rỗng để làm ranh giới mềm
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    chunks: List[str] = []
    current = ""

    for para in paragraphs:
        # Nếu thêm cả paragraph vượt quá chunk_size, flush current và bắt đầu chunk mới
        if len(current) + len(para) + 2 > chunk_size:
            if current:
                chunks.append(current.strip())
                # tạo overlap từ cuối chunk trước
                tail = current[-overlap:] if overlap > 0 else ""
                current = (tail + "\n\n" + para).strip()
            else:
                # paragraph đơn lẻ đã dài hơn chunk_size → cắt cứng theo sentence
                current = _hard_split(para, chunk_size, overlap, chunks)
                continue

            # nếu sau khi thêm mà current đã quá dài, cắt cứng tiếp
            if len(current) > chunk_size:
                current = _hard_split(current, chunk_size, overlap, chunks)
        else:
            current = (current + "\n\n" + para).strip() if current else para

    if current.strip():
        chunks.append(current.strip())

    return chunks


def _hard_split(text: str, chunk_size: int, overlap: int, out: List[str]) -> str:
    """Fallback: cắt cứng theo sentence khi một paragraph quá dài."""
    sentences = _split_sentences(text)
    buf = ""
    for sent in sentences:
        if len(buf) + len(sent) + 1 > chunk_size:
            if buf:
                out.append(buf.strip())
                tail = buf[-overlap:] if overlap > 0 else ""
                buf = (tail + " " + sent).strip()
            else:
                # Câu đơn lẻ vẫn dài hơn chunk_size → cắt theo char
                for i in range(0, len(sent), chunk_size - overlap):
                    out.append(sent[i:i + chunk_size].strip())
                buf = ""
        else:
            buf = (buf + " " + sent).strip() if buf else sent
    return buf


def _split_sentences(text: str) -> List[str]:
    """Tách câu đơn giản theo dấu chấm/chấm hỏi/chấm than/xuống dòng."""
    import re
    parts = re.split(r"(?<=[.!?])\s+|\n+", text)
    return [p.strip() for p in parts if p.strip()]


# ---------------------------------------------------------------------------
# Convenience — gọi 1 lần để parse + chunk
# ---------------------------------------------------------------------------
async def parse_and_chunk(document,
                    chunk_size: int = CHUNK_SIZE_CHARS,
                    overlap: int = CHUNK_OVERLAP_CHARS) -> tuple[List[str], List[str]]:
    """
    Helper: trích xuất text → chunk → trả về (chunks, diagrams).

    Returns:
        (chunks, diagrams): list chunk text + list mô tả diagram (có thể rỗng).
        Đây là hàm router/service khác sẽ gọi trực tiếp khi tạo Assessment.
    """
    result = await extract_text(document)
    chunks = chunk_text(result.text, chunk_size=chunk_size, overlap=overlap)
    return chunks, result.diagrams
