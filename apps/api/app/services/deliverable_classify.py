"""Stage 2 (Layer 2) — AI classify workspace files into deliverables (R1..R7/SP/SL).

Service này đóng vai trò "trọng tài AI" cho bước kiểm tra file nộp.
Nó đọc nội dung file từ MinIO, gửi AI để phân loại, rồi kết hợp với
Lớp 1 (presence check) cho kết quả cuối cùng.

Thiết kế:
- Đọc file bytes từ MinIO (qua Document.storage_key)
- Gọi AI Gateway để classify mỗi file
- Cache kết quả theo Document.content_hash (nếu có) hoặc SHA256(file_bytes)
- Fallback về Lớp 1 nếu AI lỗi/timeout

Output: dict {
    document_id ->
        {
            "deliverable_code": str | None,
            "content_ok": bool | None,
            "reason": str,
            "confidence": float (0..1)
        }
}

Mô tả issue: deliverable-check-name-only-flaw-open.md
"""

from __future__ import annotations

import base64
import hashlib
import json
import logging
import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any

from app.core.config import settings
from app.services.ai_client import ai_gateway
from app.services.storage import get_doc
from app.services.document_parser import _extract_office_images

logger = logging.getLogger(__name__)

# Cap multimodal payload: vision providers charge per image and have token limits.
# 4 images × ~1MB downscaled ≈ 4MB inline, still well under Gemini inline_data (20MB).
# For classification we only need to recognize diagram types, not pixel-perfect OCR.
_MAX_MULTIMODAL_IMAGES = 4


@dataclass
class FileClassification:
    """Kết quả classify 1 file."""
    document_id: int
    filename: str
    deliverable_code: str | None = None      # "R1".."R7", "SP", "SL", hoặc None
    content_ok: bool | None = None            # None = chưa check / lỗi, True/False
    reason: str = ""
    confidence: float = 0.0                  # 0..1


# ---------------------------------------------------------------------------
# Cache — in-memory keyed by content_hash (hoặc SHA256 nếu chưa có hash)
# ---------------------------------------------------------------------------
_classify_cache: dict[str, FileClassification] = {}


def _cache_key(document_id: int, file_bytes: bytes, content_hash: str | None) -> str:
    """Cache key: ưu tiên content_hash, fallback SHA256 file_bytes."""
    if content_hash:
        return f"classify:{content_hash}"
    return f"classify:doc_{document_id}:{hashlib.sha256(file_bytes).hexdigest()}"


# ---------------------------------------------------------------------------
# Prompt builder
# ---------------------------------------------------------------------------

# Quy tắc phân loại (chứng thực, dùng cho mọi rubric)
_CLASSIFICATION_RULES = """\
Quy tắc phân loại:
1. Chỉ trả về 1 mã deliverable phù hợp nhất (từ danh sách dưới), hoặc "unknown" nếu không xác định được.
2. Một file có thể chỉ thuộc 1 deliverable duy nhất.
3. Phân loại dựa trên NỘI DUNG file, KHÔNG dựa vào tên file.
4. Nội dung là báo cáo/markdown/word → deliverable báo cáo (R1-R7, chọn số phù hợp theo chương/nội dung).
5. Nội dung là source code/archive code → SP.
6. File rỗng, chỉ có vài dòng, hoặc nội dung không liên quan đến thesis → unknown.
7. LƯU Ý R1 vs R7: R1 = đề cương đầu kỳ (proposal); R7 = báo cáo tổng kết bảo vệ (final report/thesis) — thường nằm CUỐI file, có chương "Kết luận", "Tổng kết", "Future work". Nếu chỉ thấy đầu file (mục lục, chương 1) mà KHÔNG thấy phần kết luận → miễn cưỡng hơn khi gán R7.
"""


def _build_deliverables_list(deliverables_map: dict[str, dict[str, Any]]) -> str:
    """Build formatted deliverable list from map code->info."""
    if not deliverables_map:
        return """\
Danh sách deliverables (fallback mặc định SEP490):
- R1: Project Introduction
- R2: Project Management Plan (PMP)
- R3: Software Requirements Specifications (SRS)
- R4: Software Design Description (SDD)
- R5: Software Test Documentation (STD)
- R6: Software User Guides (SUG)
- R7: Final Project Report
- SP: Software Product (source code, .zip/.rar)
"""
    lines = ["Danh sách deliverables chính thức (từ rubric config):"]
    for code, info in sorted(deliverables_map.items()):
        name = info.get("name", code)
        file_types = info.get("file_types", [])
        ft_str = f" ({', '.join(file_types)})" if file_types else ""
        lines.append(f"- {code}: {name}{ft_str}")
    return "\n".join(lines)


def _build_classify_prompt(
    filename: str,
    text_preview: str,
    deliverables_map: dict[str, dict[str, Any]],
    max_preview_chars: int = 100000,
    total_estimated_chars: int = 50000,
) -> str:
    """Build prompt cho LLM classify file, inject deliverables từ rubric config."""
    preview = text_preview[:max_preview_chars]
    char_count = len(preview)
    ratio = round(100 * char_count / max(total_estimated_chars, 1), 1)
    if len(text_preview) > max_preview_chars:
        preview += f"\n\n[...{len(text_preview) - max_preview_chars} chars more truncated...]"
    
    deliverables_list = _build_deliverables_list(deliverables_map)
    
    return f"""\
{deliverables_list}

{_CLASSIFICATION_RULES}

=== XỬ LÝ TEXT BỊ CẮT NGẮN ===
Nội dung dưới là TRÍCH ĐOẠN bị cắt từ {char_count}/{total_estimated_chars} ký tự (tỷ lệ {ratio}%).
- Ưu tiên dùng mục lục (TOC) / tiêu đề chương để suy luận deliverable.
- Nếu KHÔNG đủ bằng chứng → confidence < 0.6 và content_ok=false, đừng đoán bừa.
- Nếu KHÔNG thấy "Kết luận/Tổng kết/Future work" mà chỉ thấy đầu file → miễn cưỡng hơn khi gán R7.

Phân loại file sau vào 1 trong các deliverables dưới đây, hoặc unknown.

Tên file: {filename}

NỘI DUNG FILE (trích đoạn):
```
{preview}
```

Trả về DUY NHẤT 1 JSON object, không có markdown, không có text khác:
{{
  "deliverable_code": "<mã từ danh sách trên>" | "unknown",
  "content_ok": true | false,
  "reason": "ngắn gọn, 1 câu, nêu bằng chứng từ nội dung",
  "confidence": 0.0
}}

Giải thích trường:
- content_ok: true = nội dung đủ/đúng chuẩn cho deliverable đó; false = rỗng, sai, hoặc thiếu phần chính/đoạn bị cắt không đủ bằng chứng.
- confidence: 0.0..1.0 mức tin cậy phân loại.
- reason: giải thích ngắn gọn tại sao chọn deliverable này.
"""


# ---------------------------------------------------------------------------
# Core classify function — gọi AI + đọc từ MinIO
# ---------------------------------------------------------------------------


async def classify_files(
    files: list[dict[str, Any]],
    deliverables: list[dict[str, Any]],
    *,
    use_cache: bool = True,
    skip_content_check: bool = False,
) -> dict[int, FileClassification]:
    """Phân loại N workspace files vào deliverables bằng AI.

    Args:
        files: list[dict] với keys: "document_id", "filename", optionally "content_hash".
               (từ workspace.files join Document)
        deliverables: list[dict] từ rubric.config["deliverables"] (để mapping code→name).
        use_cache: nếu True -> cache theo content_hash.
        skip_content_check: nếu True -> chỉ classify, không check content_ok
            (dùng khi AI bận / fallback).

    Returns:
        dict[document_id, FileClassification]
        - document_id không có trong output = AI không đọc được / lỗi
    """
    if not files:
        return {}

    result: dict[int, FileClassification] = {}
    deliverables_map = {d.get("code", ""): d for d in deliverables}
    valid_codes = set(deliverables_map.keys()) | {"unknown"}

    for f in files:
        doc_id = f.get("document_id")
        filename = f.get("filename", "unknown")
        content_hash = f.get("content_hash")

        # --- Cache check ---
        cache_k = _cache_key(doc_id, b"", content_hash) if content_hash else None
        if use_cache and cache_k and cache_k in _classify_cache:
            cached = _classify_cache[cache_k]
            result[doc_id] = FileClassification(
                document_id=doc_id,
                filename=filename,
                deliverable_code=cached.deliverable_code,
                content_ok=cached.content_ok,
                reason=f"Cached: {cached.reason}",
                confidence=cached.confidence,
            )
            continue

        # --- Read file from MinIO ---
        try:
            raw = await get_doc(f.get("storage_key", ""))  # type: ignore[arg-type]
        except Exception as exc:
            logger.warning("Cannot read file doc_id=%s from storage: %s", doc_id, exc)
            result[doc_id] = FileClassification(
                document_id=doc_id,
                filename=filename,
                reason=f"Cannot read file: {exc}",
            )
            continue

        # --- Extract text preview + images for classification ---
        text_preview, image_data_uris = _extract_text_preview(
            raw, filename, max_chars=100000
        )
        # Cap multimodal payload
        images = image_data_uris[:_MAX_MULTIMODAL_IMAGES] if image_data_uris else None

        # --- Call AI ---
        try:
            code, content_ok, reason, confidence = await _call_ai_classify(
                filename, text_preview, deliverables_map, skip_content_check,
                images=images,
            )
        except Exception as exc:
            logger.warning("AI classify failed doc_id=%s: %s", doc_id, exc)
            result[doc_id] = FileClassification(
                document_id=doc_id,
                filename=filename,
                reason=f"AI classify failed: {exc}",
            )
            continue

        classification = FileClassification(
            document_id=doc_id,
            filename=filename,
            deliverable_code=code if code in valid_codes else "unknown",
            content_ok=content_ok if not skip_content_check else None,
            reason=reason,
            confidence=confidence,
        )
        result[doc_id] = classification

        # --- Save cache ---
        if use_cache:
            if content_hash:
                _classify_cache[_cache_key(doc_id, b"", content_hash)] = classification
            else:
                _classify_cache[_cache_key(doc_id, raw, None)] = classification

    return result


# ---------------------------------------------------------------------------
# AI Call
# ---------------------------------------------------------------------------


async def _call_ai_classify(
    filename: str,
    text_preview: str,
    deliverables_map: dict[str, dict[str, Any]],
    skip_content_check: bool,
    images: list[str] | None = None,
) -> tuple[str | None, bool | None, str, float]:
    """Gọi AI để classify 1 file.

    Returns: (code, content_ok, reason, confidence)
    """
    prompt = _build_classify_prompt(filename, text_preview, deliverables_map)
    if images:
        prompt += f"\n\n[Đính kèm {len(images)} hình ảnh từ file. Hãy dùng cả text VÀ hình ảnh để phân loại.]"
    if skip_content_check:
        prompt += "\n\nLưu ý: skip_content_check=True -> đặt content_ok=null (không đánh giá nội dung)."

    resp = await ai_gateway.generate(
        prompt=prompt,
        system_prompt="Bạn là trợ lý phân loại tài liệu học thuật. Trả về JSON chính xác.",
        temperature=0.1,
        max_tokens=800,
        images=images,
    )

    # Một số provider (agnes, tencent/hy3...) tiêu token vào reasoning_content
    # (thinking) và trả content rất ngắn / bị cắt. Nếu content rỗng hoặc không
    # parse được JSON, thử lấy từ reasoning_content.
    raw_text = (resp.get("content") or "").strip()
    if not raw_text or "{" not in raw_text:
        reasoning = (
            (resp.get("raw") or {})
            .get("choices", [{}])[0]
            .get("message", {})
            .get("reasoning_content", "")
            or ""
        )
        if reasoning.strip():
            raw_text = reasoning.strip()
    parsed = _parse_classify_json(raw_text)

    if skip_content_check:
        parsed["content_ok"] = None

    return parsed.get("deliverable_code"), parsed.get("content_ok"), parsed.get("reason", ""), parsed.get("confidence", 0.0)


# ---------------------------------------------------------------------------
# Text preview extraction (lightweight, no full parse)
# ---------------------------------------------------------------------------


def _extract_text_preview(file_bytes: bytes, filename: str, max_chars: int = 100000) -> tuple[str, list[str]]:
    """Trích xuất text thô + hình ảnh (data URIs) từ file bytes.

    Returns: (text_preview, list_of_base64_data_uris)
    Giới hạn max_chars để tiết kiệm token. Mặc định 100000 để bắt gần như toàn bộ nội dung.
    """
    fname_lower = filename.lower()
    if fname_lower.endswith(".pdf"):
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(stream_bytes := file_bytes)  # type: ignore[name-defined]
            text = ""
            for page in reader.pages[:5]:  # tăng lên 5 trang đầu
                try:
                    text += (page.extract_text() or "") + "\n"
                except Exception:
                    pass
                if len(text) >= max_chars:
                    break
            return text[:max_chars], []
        except Exception as exc:
            logger.debug("PDF preview extract failed for %s: %s", filename, exc)
            return f"[PDF file, {len(file_bytes)} bytes, unable to preview text]", []

    elif fname_lower.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(BytesIO(file_bytes))
            parts = []
            # Header/footer text (metadata: tên SV, mã đồ án, ngày...)
            for section in doc.sections:
                for para in section.header.paragraphs:
                    if para.text.strip():
                        parts.append(f"[HEADER] {para.text}")
                for para in section.footer.paragraphs:
                    if para.text.strip():
                        parts.append(f"[FOOTER] {para.text}")
            # Paragraphs - tăng limit
            for para in doc.paragraphs[:200]:  # 100 -> 200
                if para.text:
                    parts.append(para.text)
            # Tables - tăng limit
            for table in doc.tables[:20]:  # 10 -> 20
                for row in table.rows[:30]:  # 20 -> 30
                    row_text = "\t".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        parts.append(row_text)

            # --- Extract embedded images as base64 data URIs ---
            image_data_uris: list[str] = []
            try:
                with zipfile.ZipFile(BytesIO(file_bytes)) as archive:
                    for info in archive.infolist():
                        if info.is_dir() or not info.filename.startswith("word/media/"):
                            continue
                        ext = info.filename.rsplit(".", 1)[-1].lower()
                        mime = {
                            "png": "image/png",
                            "jpg": "image/jpeg",
                            "jpeg": "image/jpeg",
                            "gif": "image/gif",
                            "webp": "image/webp",
                            "bmp": "image/bmp",
                        }.get(ext)
                        if not mime:
                            continue
                        raw_img = archive.read(info)
                        # Skip tiny images (icons, bullets)
                        try:
                            from PIL import Image
                            with Image.open(BytesIO(raw_img)) as img:
                                w, h = img.size
                                if w < 200 and h < 200:
                                    continue
                        except Exception:
                            pass
                        b64 = base64.b64encode(raw_img).decode()
                        image_data_uris.append(f"data:{mime};base64,{b64}")
            except Exception as exc:
                logger.debug("DOCX image extraction failed for %s: %s", filename, exc)

            return "\n".join(parts)[:max_chars], image_data_uris
        except Exception as exc:
            logger.debug("DOCX preview extract failed for %s: %s", filename, exc)
            return f"[DOCX file, {len(file_bytes)} bytes, unable to preview text]", []

    elif fname_lower.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            texts = []
            for slide in prs.slides[:5]:
                for shape in slide.shapes[:10]:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n\n".join(texts)[:max_chars], []
        except Exception as exc:
            logger.debug("PPTX preview extract failed for %s: %s", filename, exc)
            return f"[PPTX file, {len(file_bytes)} bytes, unable to preview text]", []

    elif fname_lower.endswith((".zip", ".rar")):
        return f"[Archive file {filename}, {len(file_bytes)} bytes. AI should classify as SP based on context.]", []

    return f"[File {filename}, {len(file_bytes)} bytes, type not directly readable]", []


# ---------------------------------------------------------------------------
# JSON parsing helper
# ---------------------------------------------------------------------------

_CLASSIFY_JSON_FIELDS = ("deliverable_code", "content_ok", "reason", "confidence")


def _parse_classify_json(raw: str) -> dict[str, Any]:
    """Parse JSON response from AI classify. Handle markdown fence / extra text."""
    text = raw.strip()

    # Strip markdown code fence (```json ... ``` or ``` ... ```)
    if text.startswith("```"):
        # Drop the opening fence line (``` or ```json)
        lines = text.splitlines()
        # Remove all fence lines, keep content between them
        content_lines = [ln for ln in lines if not ln.strip().startswith("```")]
        text = "\n".join(content_lines).strip()

    # Find first { and last }
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end >= start:
        text = text[start:end + 1]

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        logger.debug("Failed to parse classify JSON: %r", raw[:200])
        return {}


# ---------------------------------------------------------------------------
# Helpers for integration with Layer 1
# ---------------------------------------------------------------------------


def merge_with_presence(
    presence_result: Any,  # DeliverableCheckResult
    classifications: dict[int, FileClassification],
    deliverables: list[dict[str, Any]],
) -> Any:
    """Merge Layer 2 classifications into Layer 1 presence result.

    Logic:
    - Mỗi deliverable đươc tính "present" nếu CÓ ÍT NHẤT 1 file được AI gán vào code đó
      VÀ file đó có content_ok != False.
    - Nếu AI gán file nhưng content_ok=False → vẫn tính present=True nhưng UI hiển thị cảnh báo.
    - Nếu Layer 2 ran và produced ít nhất 1 valid classification nhưng không gán file
      vào deliverable này -> deliverable đó là "missing" (present=False), KHÔNG được
      fallback về Layer 1 (tránh false positive từ type-only matching).
    - Chỉ fallback về Layer 1 khi Layer 2 hoàn toàn thất bại (không có classification nào).

    Returns: Updated DeliverableCheckResult (with extra fields in items).
    """
    from app.services.deliverable_check import DeliverableCheckResult, DeliverableMatch

    # Determine if Layer 2 actually produced at least one valid classification.
    layer2_ok = any(
        cls.deliverable_code and cls.deliverable_code != "unknown"
        for cls in classifications.values()
    )

    # Map: deliverable_code -> list[document_id] được AI gán
    assigned: dict[str, list[int]] = {}
    for doc_id, cls in classifications.items():
        if cls.deliverable_code and cls.deliverable_code != "unknown":
            assigned.setdefault(cls.deliverable_code, []).append(doc_id)

    updated_items: list[DeliverableMatch] = []
    for item in presence_result.items:
        doc_ids_assigned = assigned.get(item.code, [])
        if doc_ids_assigned:
            # AI đã gán file -> tính present theo Layer 2
            cls = classifications[doc_ids_assigned[0]]
            filename = cls.filename
            updated_items.append(
                DeliverableMatch(
                    code=item.code,
                    name=item.name,
                    file_types=item.file_types,
                    desc=item.desc,
                    present=True,
                    matched_file=filename,
                )
            )
        elif not layer2_ok:
            # Layer 2 thất bại hoàn toàn -> giữ nguyên Layer 1
            updated_items.append(item)
        else:
            # Layer 2 chạy OK nhưng không gán file nào vào deliverable này
            # -> genuinely missing, đừng fallback về Layer 1 (avoid false positive)
            updated_items.append(
                DeliverableMatch(
                    code=item.code,
                    name=item.name,
                    file_types=item.file_types,
                    desc=item.desc,
                    present=False,
                    matched_file=None,
                )
            )

    present_count = sum(1 for it in updated_items if it.present)
    return DeliverableCheckResult(
        items=updated_items,
        total=presence_result.total,
        present_count=present_count,
    )


def clear_cache() -> None:
    """Xóa cache classify (dùng khi cần reload)."""
    _classify_cache.clear()